#!/usr/bin/env python3
"""Generate a PowerPoint presentation from the Tuas Mega Port paper."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
NAVY       = RGBColor(0x0B, 0x1D, 0x3A)
DARK_BLUE  = RGBColor(0x12, 0x2B, 0x4F)
MID_BLUE   = RGBColor(0x1A, 0x3C, 0x6E)
ACCENT     = RGBColor(0x00, 0x9E, 0xDB)
LIGHT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
GOLD       = RGBColor(0xFF, 0xB7, 0x00)
RED_ACCENT = RGBColor(0xE5, 0x3E, 0x3E)
GREEN      = RGBColor(0x4C, 0xAF, 0x50)
ORANGE     = RGBColor(0xFF, 0x98, 0x00)

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgbClr = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            alpha_elem = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_section_number(slide, number, left=Inches(0.6), top=Inches(0.4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

def add_bullet_text(text_frame, items, font_size=16, color=WHITE, spacing=Pt(8)):
    """Add bulleted items to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_before = spacing
        p.space_after = Pt(2)
        p.level = 0


TOTAL_SLIDES = 16

# =====================================================================
# SLIDE 1 – TITLE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)
# Decorative bar
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ACCENT)
add_shape_rect(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), ACCENT)
# Side accent
add_shape_rect(slide, Inches(0), Inches(1.5), Inches(0.12), Inches(4.5), ACCENT)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
             "Natural Monopoly and Systemic Risk\nin Automated Ports",
             font_size=42, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(3.8), Inches(3))
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(10), Inches(0.8),
             "An Applied Analysis of Tuas Mega Port",
             font_size=26, color=LIGHT_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(10), Inches(0.5),
             "Exploring Market Power, Structural Unemployment, and Systemic Vulnerability in Automated Hub Ports",
             font_size=15, color=LIGHT_GRAY)

# Decorative icon placeholder - port crane silhouette concept
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.8), Inches(2.0), Inches(2.8), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.color.rgb = ACCENT
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "\n\nTUAS\nMEGA\nPORT"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = 'Calibri'

add_slide_number(slide, 1, TOTAL_SLIDES)

# =====================================================================
# SLIDE 2 – ABSTRACT / OVERVIEW
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "i", Inches(0.5), Inches(0.25))
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(8), Inches(0.7),
             "Abstract & Overview", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Content cards
card1 = add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(2.8), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(1.7), Inches(3.4), Inches(0.5),
             "Research Objective", font_size=18, bold=True, color=ACCENT)
add_text_box(slide, Inches(0.7), Inches(2.2), Inches(3.4), Inches(2.0),
             "Explore excess market power in automated ports and the problems it causes: structural unemployment and vulnerability to operational/market shocks.",
             font_size=14, color=LIGHT_GRAY)

card2 = add_rounded_rect(slide, Inches(4.6), Inches(1.6), Inches(3.8), Inches(2.8), DARK_BLUE)
add_text_box(slide, Inches(4.8), Inches(1.7), Inches(3.4), Inches(0.5),
             "Key Findings", font_size=18, bold=True, color=ACCENT)
add_text_box(slide, Inches(4.8), Inches(2.2), Inches(3.4), Inches(2.0),
             "Automation boosts productivity but concentrates market power among few operators, leading to job losses and higher systemic risk.",
             font_size=14, color=LIGHT_GRAY)

card3 = add_rounded_rect(slide, Inches(8.7), Inches(1.6), Inches(4.1), Inches(2.8), DARK_BLUE)
add_text_box(slide, Inches(8.9), Inches(1.7), Inches(3.7), Inches(0.5),
             "Policy Suggestions", font_size=18, bold=True, color=ACCENT)
txBox = slide.shapes.add_textbox(Inches(8.9), Inches(2.2), Inches(3.7), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
add_bullet_text(tf, [
    "1. Labor transition programs",
    "2. Profit-sharing mechanisms",
    "3. Organization of sub-ports"
], font_size=14, color=LIGHT_GRAY)

# Outline section
add_rounded_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(4.8), Inches(4), Inches(0.5),
             "Paper Outline", font_size=20, bold=True, color=GOLD)

outline_items = [
    ("1", "Problem Definition", "Market power in automated ports"),
    ("2", "Economic Mechanisms", "Scale, lock-in, externalities"),
    ("3", "Case Study: Tuas", "Data analysis & risk assessment"),
    ("4", "Proposals", "Institutional & national solutions"),
    ("5", "Critical Review", "Assumptions, trade-offs, conclusion"),
]
for i, (num, title, desc) in enumerate(outline_items):
    x = Inches(0.7) + Inches(2.4) * i
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(5.4), Inches(0.45), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.55), Inches(5.35), Inches(1.8), Inches(0.35),
                 title, font_size=13, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.55), Inches(5.65), Inches(1.8), Inches(0.35),
                 desc, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)

# =====================================================================
# SLIDE 3 – Introduction: From Smart Cranes to Smart Monopoly
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "1")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "From Smart Cranes to Smart Monopoly", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Left content
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(7.5), Inches(1.0),
             "The New Era of Port Automation", font_size=22, bold=True, color=ACCENT)
add_text_box(slide, Inches(0.6), Inches(2.2), Inches(7.5), Inches(2.0),
             "Advancements in technology have spurred a new era of port automation:\n"
             "remotely controlled cranes, unmanned AGVs, and centralized digital\n"
             "systems cooperating with precision.",
             font_size=15, color=LIGHT_GRAY)

# Tuas highlight card
add_rounded_rect(slide, Inches(0.6), Inches(4.0), Inches(5.5), Inches(3.0), MID_BLUE)
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(5.1), Inches(0.5),
             "Tuas Mega Port at a Glance", font_size=18, bold=True, color=GOLD)

stats = [
    ("65M+ TEUs", "Target throughput capacity"),
    ("2nd Largest", "Container port globally"),
    ("1,337 ha", "Total land area"),
    ("66 Berths", "Spanning 26 km of quay"),
    ("2022", "Operations commenced"),
    ("2040", "Full completion target"),
]
for i, (val, label) in enumerate(stats):
    row = i // 2
    col = i % 2
    x = Inches(0.9) + col * Inches(2.6)
    y = Inches(4.7) + row * Inches(0.7)
    add_text_box(slide, x, y, Inches(1.2), Inches(0.35), val, font_size=16, bold=True, color=ACCENT)
    add_text_box(slide, x + Inches(1.2), y, Inches(1.5), Inches(0.35), label, font_size=12, color=LIGHT_GRAY)

# Right side - concern box
add_rounded_rect(slide, Inches(6.5), Inches(4.0), Inches(6.3), Inches(3.0), DARK_BLUE)
add_text_box(slide, Inches(6.7), Inches(4.1), Inches(5.9), Inches(0.5),
             "Concerns Raised", font_size=18, bold=True, color=RED_ACCENT)
txBox = slide.shapes.add_textbox(Inches(6.7), Inches(4.7), Inches(5.9), Inches(2.2))
tf = txBox.text_frame
tf.word_wrap = True
add_bullet_text(tf, [
    "Concentrated market power among few operators",
    "Potential monopolistic actions: price raising,\n  output restriction, market access control",
    "Reduced costs create barriers for competitors",
    "Risk of single-point-of-failure for global trade"
], font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

add_slide_number(slide, 3, TOTAL_SLIDES)

# =====================================================================
# SLIDE 4 – Key Problem & Stakeholders
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "1")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Key Problem & Stakeholder Impact", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Central problem statement
add_rounded_rect(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.3), MID_BLUE)
add_text_box(slide, Inches(1.7), Inches(1.55), Inches(9.9), Inches(1.2),
             "Core Problem: Automation at Tuas Mega Port greatly improves efficiency but creates a natural\n"
             "monopoly hub, increasing systemic risk for global trade and causing structural unemployment.",
             font_size=17, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Stakeholder cards
stakeholders = [
    ("Global Trade & Supply Chain", RED_ACCENT, [
        "Heavy tech reliance creates",
        "vulnerability to cyberattacks",
        "Severe market fluctuations",
        "from single hub disruption"
    ]),
    ("Neighboring Ports", ORANGE, [
        "Siphon Effect: customers",
        "drawn to Tuas' efficiency",
        "Matthew Effect: the strong",
        "get stronger, weak get weaker"
    ]),
    ("Port Laborers", GOLD, [
        "Structural unemployment in",
        "traditional skill-based jobs",
        "Human roles shift to monitoring",
        "from control centers"
    ]),
    ("Government", ACCENT, [
        "Must mitigate social costs via",
        "retraining programs",
        "Manage technological risks",
        "through mutual coordination"
    ]),
]

for i, (title, accent_color, points) in enumerate(stakeholders):
    x = Inches(0.5) + i * Inches(3.15)
    card = add_rounded_rect(slide, x, Inches(3.2), Inches(2.95), Inches(4.0), DARK_BLUE)
    # Accent bar at top
    add_shape_rect(slide, x, Inches(3.2), Inches(2.95), Pt(5), accent_color)
    add_text_box(slide, x + Inches(0.15), Inches(3.5), Inches(2.65), Inches(0.5),
                 title, font_size=15, bold=True, color=accent_color)
    for j, point in enumerate(points):
        add_text_box(slide, x + Inches(0.15), Inches(4.2) + j * Inches(0.5), Inches(2.65), Inches(0.5),
                     point, font_size=12, color=LIGHT_GRAY)

add_slide_number(slide, 4, TOTAL_SLIDES)

# =====================================================================
# SLIDE 5 – Economies of Scale & Market Power
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "2")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Economies of Scale & Market Power", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Left explanation
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(0.5),
             "How Natural Monopoly Forms", font_size=22, bold=True, color=ACCENT)
add_text_box(slide, Inches(0.6), Inches(2.1), Inches(6), Inches(2.5),
             "Even though the initial fixed cost of building an automated port is high,\n"
             "once built, the extra marginal cost is very low. Operational costs are\n"
             "reduced by up to 30%, creating a natural monopoly where the average\n"
             "cost for the dominant port decreases substantially in the long run.",
             font_size=15, color=LIGHT_GRAY)

# Cost comparison visual
add_rounded_rect(slide, Inches(0.6), Inches(3.8), Inches(5.8), Inches(3.2), DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(3.9), Inches(5.4), Inches(0.5),
             "Cost Structure: Automated vs Traditional Port", font_size=16, bold=True, color=GOLD)

# Visual bars showing cost breakdown
labels_vals = [
    ("Fixed Cost (Infrastructure)", 90, 40, ACCENT, RED_ACCENT),
    ("Marginal Cost per TEU", 20, 70, ACCENT, RED_ACCENT),
    ("Average Cost at Scale", 30, 75, ACCENT, RED_ACCENT),
]
for i, (label, auto_val, trad_val, auto_c, trad_c) in enumerate(labels_vals):
    y = Inches(4.6) + i * Inches(0.75)
    add_text_box(slide, Inches(0.8), y, Inches(2.5), Inches(0.3), label, font_size=11, color=LIGHT_GRAY)
    # Automated bar
    add_shape_rect(slide, Inches(3.4), y + Inches(0.02), Inches(auto_val * 0.025), Inches(0.18), auto_c)
    add_text_box(slide, Inches(3.4) + Inches(auto_val * 0.025) + Inches(0.05), y - Inches(0.02),
                 Inches(1), Inches(0.25), "Auto", font_size=9, color=ACCENT)
    # Traditional bar
    add_shape_rect(slide, Inches(3.4), y + Inches(0.25), Inches(trad_val * 0.025), Inches(0.18), trad_c)
    add_text_box(slide, Inches(3.4) + Inches(trad_val * 0.025) + Inches(0.05), y + Inches(0.22),
                 Inches(1), Inches(0.25), "Trad", font_size=9, color=RED_ACCENT)

# Right side - key insight
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5.5), MID_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.5),
             "Natural Monopoly Mechanism", font_size=20, bold=True, color=GOLD)

mechanism_steps = [
    ("1", "High Fixed Costs", "Massive infrastructure investment in automation creates\nhuge entry barriers for potential competitors"),
    ("2", "Low Marginal Costs", "Each additional container handled costs very little,\ngiving 30% operational cost reduction"),
    ("3", "Declining Average Cost", "LRATC keeps falling as quantity increases,\nmaking one large port more efficient than several"),
    ("4", "Market Dominance", "Competitors cannot match the cost structure,\nleading to natural monopoly with pricing power"),
]
for i, (num, title, desc) in enumerate(mechanism_steps):
    y = Inches(2.3) + i * Inches(1.15)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), y, Inches(0.4), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    add_text_box(slide, Inches(7.8), y - Inches(0.05), Inches(4.8), Inches(0.3),
                 title, font_size=15, bold=True, color=ACCENT)
    add_text_box(slide, Inches(7.8), y + Inches(0.25), Inches(4.8), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY)
    if i < 3:
        # Arrow down
        add_shape_rect(slide, Inches(7.37), y + Inches(0.45), Pt(3), Inches(0.5), ACCENT)

add_slide_number(slide, 5, TOTAL_SLIDES)

# =====================================================================
# SLIDE 6 – Lock-in, Switching Costs & Network Effects
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "2")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Lock-in, Switching Costs & Network Effects", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Lock-in definition
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.0), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4),
             "Lock-in Effect", font_size=20, bold=True, color=ACCENT)
add_text_box(slide, Inches(0.7), Inches(2.0), Inches(5.6), Inches(1.4),
             '"Switching away from a chosen option becomes increasingly\n'
             'difficult or expensive over time."\n'
             '— Tumisang Bogwasi',
             font_size=14, color=LIGHT_GRAY)

# Lock-in factors
add_rounded_rect(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(3.3), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(3.9), Inches(5.6), Inches(0.4),
             "Why Dominant Ports Are Hard to Replace", font_size=18, bold=True, color=GOLD)

factors = [
    "Huge number of shipping lines passing through",
    "Long-term terminal agreements in place",
    "Established relationships with labor & shipping alliances",
    "High cost of switching logistics infrastructure",
    "Network externalities amplify dominance",
]
for i, factor in enumerate(factors):
    y = Inches(4.5) + i * Inches(0.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(0.8), y, Inches(0.35), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    add_text_box(slide, Inches(1.3), y - Inches(0.02), Inches(5), Inches(0.4),
                 factor, font_size=13, color=LIGHT_GRAY)

# Matthew Effect
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5.6), MID_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4),
             "The Matthew Effect", font_size=20, bold=True, color=GOLD)
add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(1.2),
             '"Those who possess advantages tend to accumulate\n'
             'even more, while those with fewer resources fall\n'
             'further behind."',
             font_size=14, color=LIGHT_GRAY)

# Visual: two diverging arrows
add_text_box(slide, Inches(7.2), Inches(3.5), Inches(5), Inches(0.4),
             "Widening Gap Between Ports", font_size=16, bold=True, color=WHITE)

# Automated port advantages
add_rounded_rect(slide, Inches(7.2), Inches(4.1), Inches(2.4), Inches(2.8), DARK_BLUE)
add_shape_rect(slide, Inches(7.2), Inches(4.1), Inches(2.4), Pt(4), GREEN)
add_text_box(slide, Inches(7.3), Inches(4.2), Inches(2.2), Inches(0.4),
             "Automated Ports", font_size=14, bold=True, color=GREEN)
auto_items = ["More resources", "Stronger alliances", "Higher efficiency", "Pricing power"]
for i, item in enumerate(auto_items):
    add_text_box(slide, Inches(7.3), Inches(4.7) + i * Inches(0.45), Inches(2.2), Inches(0.4),
                 f"+ {item}", font_size=12, color=LIGHT_GRAY)

# Traditional port disadvantages
add_rounded_rect(slide, Inches(10.0), Inches(4.1), Inches(2.4), Inches(2.8), DARK_BLUE)
add_shape_rect(slide, Inches(10.0), Inches(4.1), Inches(2.4), Pt(4), RED_ACCENT)
add_text_box(slide, Inches(10.1), Inches(4.2), Inches(2.2), Inches(0.4),
             "Traditional Ports", font_size=14, bold=True, color=RED_ACCENT)
trad_items = ["Fewer resources", "Losing alliances", "Lower efficiency", "Price takers"]
for i, item in enumerate(trad_items):
    add_text_box(slide, Inches(10.1), Inches(4.7) + i * Inches(0.45), Inches(2.2), Inches(0.4),
                 f"- {item}", font_size=12, color=LIGHT_GRAY)

add_slide_number(slide, 6, TOTAL_SLIDES)

# =====================================================================
# SLIDE 7 – Negative Externalities
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "2")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Negative Externalities of Port Automation", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Explanation
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(6.5), Inches(2.5),
             "With little competitive pressure, operators may undermine the importance\n"
             "of safeguards and security, neglecting the socially optimal level of output\n"
             "in pursuit of profit. Extra costs are borne by other stakeholders.",
             font_size=15, color=LIGHT_GRAY)

# Three externality cards
externalities = [
    ("Cyber Risk Externality", RED_ACCENT,
     "When ports face cyberattacks,\ngovernments bear the cost of\nreducing systemic risks and\nmitigating impact on trade."),
    ("Labor Displacement", ORANGE,
     "Displaced workers suffer from\nsevere structural unemployment.\nSocial costs are externalized to\nworkers and welfare systems."),
    ("Welfare Loss", GOLD,
     "MPC shifts to MSC due to\nnegative externalities. Deadweight\nloss (area E2E1H) represents\nnet welfare loss to society."),
]
for i, (title, color, desc) in enumerate(externalities):
    x = Inches(0.5) + i * Inches(4.2)
    add_rounded_rect(slide, x, Inches(3.3), Inches(3.9), Inches(3.5), DARK_BLUE)
    add_shape_rect(slide, x, Inches(3.3), Inches(3.9), Pt(5), color)
    add_text_box(slide, x + Inches(0.2), Inches(3.6), Inches(3.5), Inches(0.4),
                 title, font_size=18, bold=True, color=color)
    add_text_box(slide, x + Inches(0.2), Inches(4.2), Inches(3.5), Inches(2.5),
                 desc, font_size=14, color=LIGHT_GRAY)

# Key takeaway
add_rounded_rect(slide, Inches(0.5), Inches(3.3 + 2.4), Inches(12.3), Inches(1.0), MID_BLUE)
add_text_box(slide, Inches(0.7), Inches(5.85), Inches(12), Inches(0.8),
             "Key Insight: The socially optimal output (Qm) is lower than the market equilibrium. "
             "The gap between MPC and MSC reflects the true social cost of concentrated port automation.",
             font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 7, TOTAL_SLIDES)

# =====================================================================
# SLIDE 8 – Case Study: Tuas Mega Port Background
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "3")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Case Study: Tuas Mega Port", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Key Assumptions
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.3), MID_BLUE)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4),
             "Key Assumptions", font_size=20, bold=True, color=GOLD)
assumptions = [
    "a. Tuas' large sunk cost keeps LRATC falling \u2192 near-natural monopoly",
    "b. No regional rival can match Tuas' automation scale",
    "c. Global container demand stays relatively stable",
]
for i, a in enumerate(assumptions):
    add_text_box(slide, Inches(0.7), Inches(2.2) + i * Inches(0.45), Inches(5.6), Inches(0.4),
                 a, font_size=13, color=LIGHT_GRAY)

# Hub vs Gateway distinction
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.3), DARK_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4),
             "Hub Port vs Gateway Port", font_size=20, bold=True, color=ACCENT)
add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(1.5),
             "Unlike gateway ports (Yangshan, Qingdao) that primarily\n"
             "receive cargo, Tuas functions as a hub port for trans-\n"
             "shipments in global supply chains. Shipping alliances\n"
             "enjoy greater flexibility in hub port selection.",
             font_size=13, color=LIGHT_GRAY)

# Port details
add_rounded_rect(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.1), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(4.2), Inches(5), Inches(0.4),
             "Tuas Mega Port Specifications", font_size=20, bold=True, color=GOLD)

specs = [
    ("Commissioning", "September 1, 2022"),
    ("Full Completion", "2040"),
    ("Total Area", "~1,337 hectares"),
    ("Berths", "66 berths spanning 26 km"),
    ("Target Capacity", "65 million+ TEUs"),
    ("Automation", "Fully automated cranes, AGVs"),
    ("Consolidation", "Absorbing Pasir Panjang, Tanjong Pagar"),
    ("Position", "World's 2nd largest container port"),
]
for i, (key, val) in enumerate(specs):
    col = i // 4
    row = i % 4
    x = Inches(0.9) + col * Inches(6)
    y = Inches(4.8) + row * Inches(0.55)
    add_text_box(slide, x, y, Inches(2), Inches(0.35), key, font_size=13, bold=True, color=ACCENT)
    add_text_box(slide, x + Inches(2.1), y, Inches(3.5), Inches(0.35), val, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 8, TOTAL_SLIDES)

# =====================================================================
# SLIDE 9 – Operational Performance Data
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "3")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Operational Performance: Before & After Tuas", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Before Tuas card
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), DARK_BLUE)
add_shape_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Pt(5), ORANGE)
add_text_box(slide, Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.4),
             "Pre-Tuas Period (2015-2022)", font_size=20, bold=True, color=ORANGE)
pre_stats = [
    ("Avg Annual Throughput", "~35.2 million TEUs"),
    ("Growth Rate", "~2.7% per year"),
    ("Range", "30.9M (2015) \u2192 37.3M (2022)"),
]
for i, (k, v) in enumerate(pre_stats):
    y = Inches(2.3) + i * Inches(0.5)
    add_text_box(slide, Inches(0.9), y, Inches(2.5), Inches(0.4), k, font_size=14, bold=True, color=LIGHT_GRAY)
    add_text_box(slide, Inches(3.5), y, Inches(2.8), Inches(0.4), v, font_size=14, color=WHITE)

# After Tuas card
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), DARK_BLUE)
add_shape_rect(slide, Inches(6.8), Inches(1.5), Inches(6), Pt(5), GREEN)
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.6), Inches(0.4),
             "Post-Tuas Period (2022-2025)", font_size=20, bold=True, color=GREEN)
post_stats = [
    ("Avg Annual Throughput", "~41.9 million TEUs"),
    ("Growth Rate", "~6.2% per year"),
    ("Range", "37.3M (2022) \u2192 44.7M (2025)"),
]
for i, (k, v) in enumerate(post_stats):
    y = Inches(2.3) + i * Inches(0.5)
    add_text_box(slide, Inches(7.2), y, Inches(2.5), Inches(0.4), k, font_size=14, bold=True, color=LIGHT_GRAY)
    add_text_box(slide, Inches(9.8), y, Inches(2.8), Inches(0.4), v, font_size=14, color=WHITE)

# Growth comparison visual
add_rounded_rect(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), MID_BLUE)
add_text_box(slide, Inches(0.7), Inches(4.4), Inches(5), Inches(0.4),
             "Throughput Growth Visualization (million TEUs)", font_size=16, bold=True, color=GOLD)

# Simple bar chart using rectangles
years_data = [
    (2015, 30.9), (2016, 31.3), (2017, 33.7), (2018, 36.6),
    (2019, 37.2), (2020, 36.9), (2021, 37.5), (2022, 37.3),
    (2023, 39.0), (2024, 42.0), (2025, 44.7),
]
bar_start_x = Inches(1.0)
bar_base_y = Inches(6.8)
bar_width = Inches(0.8)
bar_gap = Inches(0.25)
scale = Inches(0.055)  # per million TEUs

for i, (year, teu) in enumerate(years_data):
    x = bar_start_x + i * (bar_width + bar_gap)
    bar_height = teu * scale
    color = GREEN if year > 2022 else (ORANGE if year == 2022 else ACCENT)
    add_shape_rect(slide, x, bar_base_y - bar_height, bar_width, bar_height, color)
    add_text_box(slide, x - Inches(0.1), bar_base_y + Inches(0.02), Inches(1), Inches(0.3),
                 str(year), font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x - Inches(0.1), bar_base_y - bar_height - Inches(0.22), Inches(1), Inches(0.25),
                 f"{teu:.1f}", font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 9, TOTAL_SLIDES)

# =====================================================================
# SLIDE 10 – Risks of Tuas' Implementation
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "3")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Systemic Risks of Tuas' Implementation", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Turnaround time insight
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4),
             "Vessel Turnaround Time", font_size=20, bold=True, color=ACCENT)
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(1.8),
             "Average turnaround improved from 33.6 hours (2015) to ~12\n"
             "hours (2025). However, most significant gains occurred by\n"
             "2018, BEFORE Tuas. Tuas' scale expansion instead led to\n"
             "capacity overflow and congestion events.",
             font_size=14, color=LIGHT_GRAY)

# Risk Events
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), DARK_BLUE)
add_shape_rect(slide, Inches(6.8), Inches(1.5), Inches(6), Pt(5), RED_ACCENT)
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.6), Inches(0.4),
             "Congestion Crisis Events", font_size=20, bold=True, color=RED_ACCENT)

events = [
    ("2025", "Digital disruption forced ~60% of vessels\nto anchor; wait times exceeded 1 week"),
    ("2024", "Excess container discharge led to 7-day\nwaits; ~450,000 TEUs stuck at berth"),
]
for i, (year, desc) in enumerate(events):
    y = Inches(2.3) + i * Inches(0.85)
    add_text_box(slide, Inches(7.2), y, Inches(0.8), Inches(0.3), year, font_size=16, bold=True, color=GOLD)
    add_text_box(slide, Inches(8.1), y, Inches(4.5), Inches(0.8), desc, font_size=12, color=LIGHT_GRAY)

# Employment impact
add_rounded_rect(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), MID_BLUE)
add_text_box(slide, Inches(0.7), Inches(4.4), Inches(5.6), Inches(0.4),
             "Structural Unemployment Impact", font_size=20, bold=True, color=GOLD)

add_text_box(slide, Inches(0.7), Inches(5.0), Inches(5.5), Inches(2.0),
             "Tuas is equipped with automated cranes and driverless\n"
             "AGVs. Humans now monitor operations from control centers\n"
             "rather than performing skill-based work, spurring job\n"
             "losses in traditional port positions.",
             font_size=14, color=LIGHT_GRAY)

# Impact visual
impact_items = [
    ("Traditional Crane Operators", "Replaced", RED_ACCENT),
    ("Manual Cargo Handlers", "Replaced", RED_ACCENT),
    ("Vehicle / Truck Drivers", "Replaced by AGVs", RED_ACCENT),
    ("Control Center Monitors", "New Role", GREEN),
    ("IT / System Engineers", "New Role", GREEN),
    ("Data Analysts", "New Role", GREEN),
]
for i, (role, status, color) in enumerate(impact_items):
    col = i // 3
    row = i % 3
    x = Inches(6.8) + col * Inches(3)
    y = Inches(4.9) + row * Inches(0.6)
    add_text_box(slide, x, y, Inches(2.2), Inches(0.3), role, font_size=12, color=LIGHT_GRAY)
    add_text_box(slide, x + Inches(2.2), y, Inches(1), Inches(0.3), status, font_size=12, bold=True, color=color)

add_slide_number(slide, 10, TOTAL_SLIDES)

# =====================================================================
# SLIDE 11 – Institutional Proposal: Labor Transition
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "4")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Proposal 1: Labor Transition & Profit-Sharing", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Phase 1
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(3.9), Inches(4.5), DARK_BLUE)
add_shape_rect(slide, Inches(0.5), Inches(1.5), Inches(3.9), Pt(5), ACCENT)
add_text_box(slide, Inches(0.7), Inches(1.7), Inches(3.5), Inches(0.4),
             "Phase 1: Skills & Partnerships", font_size=17, bold=True, color=ACCENT)
phase1 = [
    "Identify skill gaps in workforce",
    "Partner with local vocational colleges",
    "Free training in: AI scheduling, data\n  recording, equipment maintenance,\n  in-time debugging",
    "Skill certificates recognized by\n  other ports",
    "Guaranteed job offer upon completion",
    "10-15% salary increase for trainees",
]
for i, item in enumerate(phase1):
    add_text_box(slide, Inches(0.7), Inches(2.3) + i * Inches(0.6), Inches(3.5), Inches(0.55),
                 f"\u2022 {item}", font_size=11, color=LIGHT_GRAY)

# Phase 2
add_rounded_rect(slide, Inches(4.7), Inches(1.5), Inches(3.9), Inches(4.5), DARK_BLUE)
add_shape_rect(slide, Inches(4.7), Inches(1.5), Inches(3.9), Pt(5), GOLD)
add_text_box(slide, Inches(4.9), Inches(1.7), Inches(3.5), Inches(0.4),
             "Phase 2: Smooth Transitions", font_size=17, bold=True, color=GOLD)
phase2 = [
    "Priority internal transfers:\n  automated roles first offered\n  to existing workers",
    "No external competition for\n  initial placement",
    "6-month transition period with\n  original salary guaranteed",
    "Even if new role pays less\n  initially, salary protected",
]
for i, item in enumerate(phase2):
    add_text_box(slide, Inches(4.9), Inches(2.3) + i * Inches(0.75), Inches(3.5), Inches(0.7),
                 f"\u2022 {item}", font_size=11, color=LIGHT_GRAY)

# Phase 3
add_rounded_rect(slide, Inches(8.9), Inches(1.5), Inches(3.9), Inches(4.5), DARK_BLUE)
add_shape_rect(slide, Inches(8.9), Inches(1.5), Inches(3.9), Pt(5), GREEN)
add_text_box(slide, Inches(9.1), Inches(1.7), Inches(3.5), Inches(0.4),
             "Phase 3: Profit & Motivation", font_size=17, bold=True, color=GREEN)
phase3 = [
    "20% of total revenue allocated:",
    "  \u2022 5% for savings (partnerships)",
    "  \u2022 15% for worker bonuses",
    "Higher skill = bigger bonus",
    "Workers motivated to upskill",
    "Creates long-term retention",
]
for i, item in enumerate(phase3):
    add_text_box(slide, Inches(9.1), Inches(2.3) + i * Inches(0.6), Inches(3.5), Inches(0.55),
                 f"{item}", font_size=11, color=LIGHT_GRAY)

# Mutual benefit box
add_rounded_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1.0), MID_BLUE)
add_text_box(slide, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.8),
             "Reciprocity: Colleges gain access to advanced devices & field practice  |  "
             "Ports enjoy education resources  |  Workers gain transferable skills & job security",
             font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 11, TOTAL_SLIDES)

# =====================================================================
# SLIDE 12 – Evaluation of Institutional Proposal
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "4")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Evaluation: Institutional Proposal", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Pros
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(8.5), Inches(5.2), DARK_BLUE)
add_shape_rect(slide, Inches(0.5), Inches(1.5), Inches(8.5), Pt(5), GREEN)
add_text_box(slide, Inches(0.7), Inches(1.7), Inches(4), Inches(0.4),
             "Advantages", font_size=22, bold=True, color=GREEN)

pros = [
    ("Less Unemployment", "No sudden job loss for internal workers;\nguaranteed placement + better pay through bonuses;\ntransferable skills for any automated port."),
    ("Minimized Market Power", "Skilled workers across all ports reduce the\nMatthew and Siphon effects; every port gains\ncapability, not just the dominant one."),
    ("Competitive Supply Chain", "Diminished market power \u2192 competition \u2192\nmore affordable pricing, abundant & stable\nsupply, and higher quality goods."),
    ("Mutual Benefit", "Colleges: advanced devices & field practice.\nPorts: education resources.\nWorkers: better knowledge & security."),
]
for i, (title, desc) in enumerate(pros):
    y = Inches(2.3) + i * Inches(1.05)
    add_text_box(slide, Inches(0.9), y, Inches(2.2), Inches(0.3), title, font_size=14, bold=True, color=ACCENT)
    add_text_box(slide, Inches(3.2), y, Inches(5.5), Inches(0.9), desc, font_size=12, color=LIGHT_GRAY)

# Limitations
add_rounded_rect(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(5.2), DARK_BLUE)
add_shape_rect(slide, Inches(9.3), Inches(1.5), Inches(3.5), Pt(5), RED_ACCENT)
add_text_box(slide, Inches(9.5), Inches(1.7), Inches(3.1), Inches(0.4),
             "Limitations", font_size=22, bold=True, color=RED_ACCENT)
add_text_box(slide, Inches(9.5), Inches(2.3), Inches(3.1), Inches(4.2),
             "Increases short-term\noperational costs:\n\n"
             "\u2022 Training investment\n"
             "\u2022 Career coaching\n"
             "\u2022 Role restructuring\n\n"
             "However, long-run benefits\noutweigh short-term costs,\nmaking the proposal\nmostly advantageous.",
             font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 12, TOTAL_SLIDES)

# =====================================================================
# SLIDE 13 – National Proposal: Sub-Ports
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "4")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Proposal 2: Re-Balancing via Sub-Ports", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Strategy description
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.8),
             "Organize main-ports and sub-ports to relieve shipping concentration and vulnerability\n"
             "while preserving the natural monopoly's efficiency advantages.",
             font_size=16, color=LIGHT_GRAY)

# The hub-and-spoke diagram concept
add_rounded_rect(slide, Inches(0.5), Inches(2.5), Inches(6), Inches(4.7), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(2.6), Inches(5.6), Inches(0.4),
             "Hub & Sub-Port Model", font_size=20, bold=True, color=GOLD)

# Central hub
hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), Inches(3.8), Inches(2), Inches(1.5))
hub.fill.solid()
hub.fill.fore_color.rgb = ACCENT
hub.line.fill.background()
tf = hub.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "TUAS\nMEGA PORT\n(Main Hub)"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Sub-ports
sub_positions = [
    (Inches(0.8), Inches(3.2), "Pasir\nPanjang"),
    (Inches(4.5), Inches(3.0), "Sub-Port\nB"),
    (Inches(0.8), Inches(5.5), "Sub-Port\nC"),
    (Inches(4.5), Inches(5.7), "Sub-Port\nD"),
]
for x, y, name in sub_positions:
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(1.2), Inches(0.9))
    sp.fill.solid()
    sp.fill.fore_color.rgb = MID_BLUE
    sp.line.color.rgb = ACCENT
    sp.line.width = Pt(1.5)
    tf = sp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

# Right side - Policy details
add_rounded_rect(slide, Inches(6.8), Inches(2.5), Inches(6), Inches(4.7), MID_BLUE)
add_text_box(slide, Inches(7.0), Inches(2.6), Inches(5.6), Inches(0.4),
             "Implementation Strategy", font_size=20, bold=True, color=GOLD)

strategies = [
    ("Government Subsidies", "Provide funding for sub-port infrastructure.\nSingapore's 0.9% budget surplus (2025) and\n$857B GDP demonstrate fiscal capacity."),
    ("Retain Pasir Panjang", "Instead of phasing out when Tuas reaches\nfull capacity, keep as active sub-port to\nmaintain distributed shipping capacity."),
    ("PSA Coordination", "Both Tuas and sub-ports operated under PSA,\nenabling centralized but distributed management\nof Singapore's port ecosystem."),
    ("Traffic Distribution", "Redistribute shipping lines across sub-ports to\nreduce single-point-of-failure risk and alleviate\ncongestion at the main hub."),
]
for i, (title, desc) in enumerate(strategies):
    y = Inches(3.1) + i * Inches(1.0)
    add_text_box(slide, Inches(7.2), y, Inches(5.4), Inches(0.3),
                 title, font_size=14, bold=True, color=ACCENT)
    add_text_box(slide, Inches(7.2), y + Inches(0.3), Inches(5.4), Inches(0.7),
                 desc, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 13, TOTAL_SLIDES)

# =====================================================================
# SLIDE 14 – Evaluation of National Proposal
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "4")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Evaluation: National Proposal", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Pros section
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(8.5), Inches(5.5), DARK_BLUE)
add_shape_rect(slide, Inches(0.5), Inches(1.5), Inches(8.5), Pt(5), GREEN)
add_text_box(slide, Inches(0.7), Inches(1.7), Inches(4), Inches(0.4),
             "Advantages", font_size=22, bold=True, color=GREEN)

nat_pros = [
    ("Positive Feedback Loop",
     "Government investment \u2192 infrastructure & jobs at sub-ports\n"
     "\u2192 improved logistics \u2192 more competitive industries\n"
     "\u2192 increased trade & tax revenue \u2192 further investment"),
    ("Creates Competition",
     "Scattered port distribution creates genuine competition,\n"
     "incentivizing Tuas to keep improving technology\n"
     "while providing alternatives during disruptions."),
    ("Relieves Unemployment",
     "Combined with institutional proposal, workers can\n"
     "transition to not-yet-fully-automated sub-ports,\n"
     "mitigating displacement from automation."),
    ("Supply Chain Resilience",
     "Multiple operational ports reduce the risk of\n"
     "supply chain breakdown when technical problems\n"
     "affect a single hub."),
]
for i, (title, desc) in enumerate(nat_pros):
    y = Inches(2.3) + i * Inches(1.15)
    add_text_box(slide, Inches(0.9), y, Inches(2.5), Inches(0.3), title, font_size=14, bold=True, color=ACCENT)
    add_text_box(slide, Inches(3.5), y, Inches(5.3), Inches(1.0), desc, font_size=12, color=LIGHT_GRAY)

# Limitations
add_rounded_rect(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(5.5), DARK_BLUE)
add_shape_rect(slide, Inches(9.3), Inches(1.5), Inches(3.5), Pt(5), RED_ACCENT)
add_text_box(slide, Inches(9.5), Inches(1.7), Inches(3.1), Inches(0.4),
             "Limitations", font_size=22, bold=True, color=RED_ACCENT)
add_text_box(slide, Inches(9.5), Inches(2.3), Inches(3.1), Inches(4.5),
             "Resource Allocation\nInefficiency:\n\n"
             "Tuas as natural monopolist\nhas extremely low marginal\ncost per unit of service.\n\n"
             "Adding sub-ports reduces\nTuas throughput, under-\nutilizing its cost advantage.\n\n"
             "Sub-ports operate at\nhigher costs, making the\noverall allocation\ncomparatively less efficient.",
             font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 14, TOTAL_SLIDES)

# =====================================================================
# SLIDE 15 – Critical Review & Trade-offs
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_section_number(slide, "5")
add_text_box(slide, Inches(1.4), Inches(0.25), Inches(10), Inches(0.7),
             "Critical Review: Assumptions & Trade-offs", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(1.4), Inches(1.05), Inches(2.5))

# Key Assumptions
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.8), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4),
             "Key Assumptions", font_size=20, bold=True, color=GOLD)
assumptions_list = [
    "Automated ports work well enough to achieve market power",
    "Labor faces structural unemployment without ability\nto improve skills or cooperate with technology",
    "Government is willing to coordinate ports and\nshipping lines",
    "Government is willing and able to fund interventions",
]
for i, a in enumerate(assumptions_list):
    add_text_box(slide, Inches(0.9), Inches(2.2) + i * Inches(0.55), Inches(5.4), Inches(0.5),
                 f"{i+1}. {a}", font_size=13, color=LIGHT_GRAY)

# Trade-offs
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.8), MID_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4),
             "Key Trade-offs", font_size=20, bold=True, color=GOLD)

tradeoffs = [
    ("Efficiency vs. Resilience", "Automation boosts capacity but creates\nsingle-point vulnerabilities"),
    ("Cost Savings vs. Employment", "Lower operational costs come at the\nprice of structural job losses"),
    ("Monopoly vs. Competition", "Natural monopoly is efficient but\nreduces market competition"),
]
for i, (title, desc) in enumerate(tradeoffs):
    y = Inches(2.2) + i * Inches(0.65)
    add_text_box(slide, Inches(7.2), y, Inches(2.5), Inches(0.3), title, font_size=13, bold=True, color=ACCENT)
    add_text_box(slide, Inches(9.8), y, Inches(2.8), Inches(0.55), desc, font_size=11, color=LIGHT_GRAY)

# Policy Balance
add_rounded_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), DARK_BLUE)
add_text_box(slide, Inches(0.7), Inches(4.7), Inches(5.6), Inches(0.4),
             "Balancing the Trade-offs", font_size=20, bold=True, color=GOLD)

balance_items = [
    ("Institutional Level", "Training + profit-sharing mitigates unemployment\nwhile building workforce capacity across all ports"),
    ("National Level", "Sub-port organization preserves natural monopoly\nefficiency while distributing systemic risk"),
    ("Combined Effect", "Short-term costs are insignificant compared to\nlong-term stability, competition, and resilience"),
]
for i, (title, desc) in enumerate(balance_items):
    x = Inches(0.8) + i * Inches(4.1)
    add_text_box(slide, x, Inches(5.2), Inches(3.8), Inches(0.3),
                 title, font_size=15, bold=True, color=ACCENT)
    add_text_box(slide, x, Inches(5.6), Inches(3.8), Inches(1.2),
                 desc, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 15, TOTAL_SLIDES)

# =====================================================================
# SLIDE 16 – Conclusion
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ACCENT)
add_shape_rect(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), ACCENT)
add_shape_rect(slide, Inches(0), Inches(1.5), Inches(0.12), Inches(4.5), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8),
             "Conclusion", font_size=38, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.3), Inches(3))

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
             "The solutions proposed — training workers for new roles, supporting job transitions,\n"
             "and spreading port operations across smaller sites — can balance the trade-offs\n"
             "between automation efficiency and its negative consequences.",
             font_size=18, color=LIGHT_GRAY)

# Three conclusion pillars
pillars = [
    ("Competitive Ports", ACCENT,
     "Keep Singapore's ports\nglobally competitive through\nautomation efficiency"),
    ("Protected Workers", GREEN,
     "Shield workers from displacement\nthrough training, profit-sharing,\nand smooth transitions"),
    ("Stable System", GOLD,
     "Make the maritime system\nmore resilient through\ndistributed operations"),
]
for i, (title, color, desc) in enumerate(pillars):
    x = Inches(0.8) + i * Inches(4.1)
    add_rounded_rect(slide, x, Inches(3.8), Inches(3.8), Inches(2.2), DARK_BLUE)
    add_shape_rect(slide, x, Inches(3.8), Inches(3.8), Pt(5), color)
    add_text_box(slide, x + Inches(0.2), Inches(4.0), Inches(3.4), Inches(0.4),
                 title, font_size=18, bold=True, color=color)
    add_text_box(slide, x + Inches(0.2), Inches(4.5), Inches(3.4), Inches(1.5),
                 desc, font_size=14, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7),
             "By following these proposals, Singapore can maintain port competitiveness while\n"
             "protecting workers and ensuring long-term maritime system stability.",
             font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 16, TOTAL_SLIDES)

# ── Save ──
output_path = "/home/user/scc2/Tuas_Mega_Port_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
